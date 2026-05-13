"""Extract PIPG PPT (Pondok Indah 골프장 운영 분석.pptx) into structured JSON.
Output: site/data/pipg_pptx_data.json"""
import sys, json, os, re
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation

PPT_PATH = 'C:/Users/yoonseok.moon/OneDrive - (주) ST International/Projects/Matoa 골프장/Pondok Indah 골프장 운영 분석.pptx'
OUT = 'C:/Users/yoonseok.moon/OneDrive - (주) ST International/Projects/Matoa 골프장/site/data/pipg_pptx_data.json'

p = Presentation(PPT_PATH)

def get_table(slide_idx, table_idx=0):
    slide = p.slides[slide_idx-1]
    tables = [sh.table for sh in slide.shapes if sh.has_table]
    if table_idx >= len(tables): return None
    t = tables[table_idx]
    rows = []
    for row in t.rows:
        rows.append([cell.text_frame.text.strip() for cell in row.cells])
    return rows

def to_num(s):
    if not s or s == '': return None
    s = s.replace(',','').replace(' ','').strip()
    if s.startswith('(') and s.endswith(')'):
        try: return -int(s[1:-1])
        except: return None
    if s == '-' or s == '—' or s == 'N/A' or s == '%': return None
    try: return int(s)
    except:
        try: return float(s)
        except: return None

def to_pct(s):
    if not s: return None
    s = s.replace('%','').replace(' ','').strip()
    if s.startswith('(') and s.endswith(')'):
        try: return -float(s[1:-1])
        except: return None
    try: return float(s)
    except: return None

out = {
    'source': 'Pondok Indah 골프장 운영 분석.pptx (Slide 1-13, internal slide deck dated 2025-11-13)',
    'extracted': '2026-05-13',
    'unit': 'IDR Million unless noted',
    'fiscal_years': [str(y) for y in range(2011, 2025)],  # FY2011-FY2024
}

# === Slide 2: Land rights inventory (HGB/HP) ===
t = get_table(2, 0)
if t:
    items = []
    total = 0
    for row in t[1:]:
        status = row[0].strip()
        if not status or '합계' in status.lower() or 'total' in status.lower() or 'large' in status.lower():
            continue
        validity = row[1].strip()
        try: area = int(row[2].replace(',','').replace(' ','').strip())
        except: area = None
        if status and not status.lower().startswith('합'):
            items.append({'parcel': status, 'validity': validity, 'area_m2': area})
            if area: total += area
    out['land_rights'] = {
        'items': items,
        'total_m2': total,
        'count': len(items),
        'note': 'HGB/HP from PIPG AR Note 10. Cap. 530,095 m² confirmed in deck.'
    }

# === Slide 3: Balance sheet 14 years (41 rows × 15 cols) ===
t = get_table(3, 0)
if t and len(t) > 1:
    years = [c.strip() for c in t[0][1:]]
    bs = {}
    for row in t[1:]:
        label = row[0].strip()
        if not label: continue
        values = [to_num(c) for c in row[1:1+len(years)]]
        bs[label] = dict(zip(years, values))
    out['balance_sheet'] = {'years': years, 'lines': bs}

# === Slide 4: P&L 14 years ===
t = get_table(4, 0)
if t and len(t) > 1:
    years = [c.strip() for c in t[0][1:]]
    pl = {}
    for row in t[1:]:
        label = row[0].strip()
        if not label: continue
        values = [to_num(c) for c in row[1:1+len(years)]]
        pl[label] = dict(zip(years, values))
    out['pnl'] = {'years': years, 'lines': pl}

# === Slide 6: Traffic — 3 tables ===
# Table 1: Visitor counts (golf course)
t = get_table(6, 0)
if t:
    visitors = []
    for row in t[3:]:  # header is rows 0-2
        try:
            year = int(row[0].strip())
            members = to_num(row[1])
            non_members = to_num(row[2])
            total = to_num(row[3])
            daily_365 = to_num(row[4])
            daily_330 = to_num(row[5])
            teams_365 = to_num(row[6])
            teams_330 = to_num(row[7])
            visitors.append({
                'year': year, 'members': members, 'non_members': non_members,
                'total': total, 'avg_daily_365': daily_365, 'avg_daily_330': daily_330,
                'avg_teams_365': teams_365, 'avg_teams_330': teams_330,
            })
        except: pass
    out['traffic_golf'] = visitors

# Table 2: Driving Range
t = get_table(6, 1)
if t:
    dr = []
    for row in t[3:]:
        try:
            year = int(row[0].strip())
            members = to_num(row[1])
            non_members = to_num(row[2])
            total = to_num(row[3])
            daily = to_num(row[4])
            dr.append({'year': year, 'members': members, 'non_members': non_members, 'total': total, 'avg_daily': daily})
        except: pass
    out['traffic_driving_range'] = dr

# Table 3: Per-capita spend (14 rows, with year col)
t = get_table(6, 2)
if t:
    spend = []
    for row in t[2:]:
        try:
            year = int(row[0].strip())
            spend.append({
                'year': year,
                'golf_course_idr_mil': to_num(row[1]),
                'cart_idr_mil': to_num(row[2]),
                'golf_total_idr_mil': to_num(row[3]),
                'golf_per_visitor_idr': to_num(row[4]),
                'restaurant_idr_mil': to_num(row[5]),
                'restaurant_per_visitor_idr': to_num(row[6]),
                'total_idr_mil': to_num(row[7]),
                'total_per_visitor_idr': to_num(row[8]),
            })
        except: pass
    out['per_capita_spend'] = spend

# === Slide 7: Sector revenue / COGS / Gross profit (3 tables × 12 segments × 14 years) ===
def extract_segment_table(table_idx, key):
    t = get_table(7, table_idx)
    if not t: return None
    years_row = t[0]
    years = [c.strip() for c in years_row[1:15] if c.strip()]
    segments = {}
    for row in t[1:]:
        label = row[0].strip()
        if not label: continue
        # First 14 cols after label are years, last col is share %
        values = [to_num(c) for c in row[1:15]]
        share = row[15] if len(row) > 15 else ''
        segments[label] = {'yearly': dict(zip(years, values)), 'share_label': share.strip()}
    return {'years': years, 'segments': segments}

out['segment_revenue']    = extract_segment_table(0, 'rev')
out['segment_cogs']       = extract_segment_table(1, 'cogs')
out['segment_gross_profit'] = extract_segment_table(2, 'gp')

# Save
with open(OUT, 'w', encoding='utf-8') as f:
    json.dump(out, f, ensure_ascii=False, indent=2)

# Summary
print(f'Wrote {OUT}')
print(f'  Land rights: {len(out.get("land_rights",{}).get("items",[]))} records, total {out.get("land_rights",{}).get("total_m2",0):,} m²')
print(f'  P&L lines: {len(out.get("pnl",{}).get("lines",{}))} × {len(out.get("pnl",{}).get("years",[]))} years')
print(f'  Balance sheet lines: {len(out.get("balance_sheet",{}).get("lines",{}))}')
print(f'  Segment revenue: {len(out.get("segment_revenue",{}).get("segments",{}))} segments')
print(f'  Segment COGS: {len(out.get("segment_cogs",{}).get("segments",{}))} segments')
print(f'  Traffic golf: {len(out.get("traffic_golf",[]))} years')
print(f'  Traffic DR: {len(out.get("traffic_driving_range",[]))} years')
print(f'  Per-capita spend: {len(out.get("per_capita_spend",[]))} years')
