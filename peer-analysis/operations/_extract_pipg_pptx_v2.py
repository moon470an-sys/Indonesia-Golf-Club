"""PIPG PPT extraction v2 - adds 5 sections to existing pipg_pptx_data.json:
  capex_ppe_14y      (Slide 8)  : 13 PPE category × 14y (new_invest, depreciation, book_value)
  capex_auxiliary    (Slide 9)  : investment property + deferred land rights + RoU
  opex_14y           (Slide 10) : 18 OPEX line items × 14y + CAGR + averages
  organization_13y   (Slide 11) : headcount (12 dept × 13y) + labor cost
  pbb_13y            (Slide 12) : land tax (PBB) area + assessed value + tax × 13y
  contracts          (Slide 13) : 21 commercial contracts + water usage

Strategy: read existing pipg_pptx_data.json, append/update the new sections, write back.
"""
import sys, json, os, re
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation

PPT_PATH = r'C:\Users\yoonseok.moon\OneDrive - (주) ST International\Projects\Matoa 골프장\Pondok Indah 골프장 운영 분석.pptx'
OUT      = r'C:\Users\yoonseok.moon\OneDrive - (주) ST International\Projects\Matoa 골프장\site\data\pipg_pptx_data.json'

p = Presentation(PPT_PATH)

with open(OUT, encoding='utf-8') as f:
    out = json.load(f)

def cells(slide_idx, tbl_idx=0):
    slide = p.slides[slide_idx-1]
    tables = [sh.table for sh in slide.shapes if sh.has_table]
    if tbl_idx >= len(tables): return None
    t = tables[tbl_idx]
    return [[c.text_frame.text.strip() for c in r.cells] for r in t.rows]

def to_num(s):
    if s is None: return None
    s = s.strip().replace(',', '').replace(' ', '')
    if not s or s in ('-', '—', 'N/A', '%'): return None
    neg = False
    if s.startswith('(') and s.endswith(')'):
        s = s[1:-1]; neg = True
    s = s.rstrip('%')
    try:
        v = int(s) if '.' not in s else float(s)
        return -v if neg else v
    except: return None

def to_pct(s):
    if not s: return None
    m = re.search(r'\(?(-?[\d.]+)%?\)?', s)
    if not m: return None
    try: return float(m.group(1))
    except: return None

# =====================================================================
# Slide 8: CAPEX 분석 - 유형자산 (PPE)
# Table 0: label | 2010 BV | (2011 NewInv, Dep, BV) | ... | (2017 NewInv, Dep, BV)  -> 1 + 1 + 7*3 = 23 cols
# Table 1: label | (2018 NewInv, Dep, BV) | ... | (2024 NewInv, Dep, BV) | 연평균(3) -> 1 + 7*3 + 3 = 25 cols
# Rows: header(2) + 13 categories incl. Total -> 15 rows total
# =====================================================================
def parse_capex_table(t, label_col, year_start_col, years_layout):
    """years_layout: list of (year_label, n_cells) tuples
    For 2010 it's 1 cell (BV only); others 3 cells (NewInv, Dep, BV).
    """
    categories = {}
    # data rows start at index 2 (rows 0-1 are headers)
    for row in t[2:]:
        label = row[label_col].strip()
        if not label: continue
        idx = year_start_col
        yearly = {}
        for ylabel, n in years_layout:
            if n == 1:
                yearly[ylabel] = {'book_value': to_num(row[idx]) if idx < len(row) else None}
            else:
                yearly[ylabel] = {
                    'new_investment': to_num(row[idx]) if idx < len(row) else None,
                    'depreciation_etc': to_num(row[idx+1]) if idx+1 < len(row) else None,
                    'book_value': to_num(row[idx+2]) if idx+2 < len(row) else None,
                }
            idx += n
        categories[label] = yearly
    return categories

t0 = cells(8, 0)  # FY10-FY17
t1 = cells(8, 1)  # FY18-FY24 + 연평균
if t0 and t1:
    capex_ppe = {}
    # Table 0 layout: 2010 (1 cell) + 2011..2017 (3 cells each)
    layout0 = [('2010', 1)] + [(str(y), 3) for y in range(2011, 2018)]
    cats0 = parse_capex_table(t0, 0, 1, layout0)
    # Table 1 layout: 2018..2024 (3 cells each) + 연평균 (3 cells)
    layout1 = [(str(y), 3) for y in range(2018, 2025)] + [('avg_11_25', 3)]
    cats1 = parse_capex_table(t1, 0, 1, layout1)
    # Merge by category
    all_cats = sorted(set(cats0.keys()) | set(cats1.keys()))
    for cat in all_cats:
        merged = {}
        if cat in cats0: merged.update(cats0[cat])
        if cat in cats1: merged.update(cats1[cat])
        capex_ppe[cat] = merged
    out['capex_ppe_14y'] = {
        'source_slide': 8,
        'source_title': '7. CAPEX 분석 - 유형자산',
        'years': [str(y) for y in range(2010, 2025)] + ['avg_11_25'],
        'unit': 'IDR Million',
        'columns_per_year': {'2010': ['book_value'], 'other': ['new_investment', 'depreciation_etc', 'book_value']},
        'categories': capex_ppe,
        'caveat': 'avg_11_25 column is the source\'s "연평균(\'11~\'25)" footer; cell alignment is inconsistent across rows so its values may not be reliable. Derive averages client-side from yearly book_value/new_investment instead.'
    }

# =====================================================================
# Slide 9: CAPEX 보조 (투자부동산, 이연토지사용권, 사용권자산)
# Table 0+1: 투자부동산 split FY10-17 / FY18-24
# Table 2: 이연토지사용권 (HGB/HP) FY11-24
# Table 3: 사용권자산 (RoU MKPI) FY23-24
# Same parser as PPE for tables 0/1.
# =====================================================================
t0 = cells(9, 0)
t1 = cells(9, 1)
t2 = cells(9, 2)
t3 = cells(9, 3)
aux = {}
if t0 and t1:
    layout0 = [('2010', 1)] + [(str(y), 3) for y in range(2011, 2018)]
    layout1 = [(str(y), 3) for y in range(2018, 2025)]
    a = parse_capex_table(t0, 0, 1, layout0)
    b = parse_capex_table(t1, 0, 1, layout1)
    inv_prop = {}
    for cat in sorted(set(a.keys()) | set(b.keys())):
        merged = {}
        if cat in a: merged.update(a[cat])
        if cat in b: merged.update(b[cat])
        inv_prop[cat] = merged
    aux['investment_property'] = inv_prop
# Deferred land rights (table 2) is messy header — store raw cells for now
if t2:
    aux['deferred_land_rights_raw'] = t2
if t3:
    aux['rou_asset_raw'] = t3
if aux:
    out['capex_auxiliary'] = {
        'source_slide': 9,
        'source_title': '8. CAPEX 분석 – 투자 부동산 / 사용권자산 / 이연토지사용권',
        'unit': 'IDR Million',
        **aux
    }

# =====================================================================
# Slide 10: OPEX 14년
# Table 0: 20 rows × 19 cols
# Header[0]:  ['(단위)', '2011', '2012', ..., '2024', '', 'CAGR', '연평균', '']
# Row layout: 18 OPEX lines + Total + ?
# Cols:        label | 14y values | share%FY24 | CAGR | avg | avg_share%
# =====================================================================
t = cells(10, 0)
if t:
    header = t[0]
    # Identify year columns
    year_cols = []
    for i, h in enumerate(header):
        if re.match(r'^20\d{2}$', h.strip()):
            year_cols.append((i, h.strip()))
    years = [y for _, y in year_cols]
    cagr_col = None
    avg_col = None
    share_fy24_col = None
    avg_share_col = None
    for i, h in enumerate(header):
        if 'CAGR' in h: cagr_col = i
        elif '연평균' in h: avg_col = i
    if year_cols:
        last_year_idx = year_cols[-1][0]
        # share% FY24 = first non-empty cell after last year column
        share_fy24_col = last_year_idx + 1
    # avg_share is the column after avg_col
    if avg_col is not None:
        avg_share_col = avg_col + 1

    opex_lines = []
    total_row = None
    for row in t[1:]:
        label = row[0].strip()
        if not label: continue
        rec = {
            'label_ko': label,
            'yearly': {y: to_num(row[i]) for i, y in year_cols if i < len(row)},
        }
        if share_fy24_col is not None and share_fy24_col < len(row):
            rec['share_fy24_pct'] = to_pct(row[share_fy24_col])
        if cagr_col is not None and cagr_col < len(row):
            rec['cagr_11_24_pct'] = to_pct(row[cagr_col])
        if avg_col is not None and avg_col < len(row):
            rec['avg_11_24'] = to_num(row[avg_col])
        if avg_share_col is not None and avg_share_col < len(row):
            rec['avg_share_pct'] = to_pct(row[avg_share_col])
        if label == 'Total':
            total_row = rec
        else:
            opex_lines.append(rec)
    out['opex_14y'] = {
        'source_slide': 10,
        'source_title': '9. OPEX 분석',
        'years': years,
        'unit': 'IDR Million',
        'lines': opex_lines,
        'total': total_row,
        'note': '18 OPEX 항목 + Total. share/avg/CAGR 칼럼 별도 추출. PBB는 별도 슬라이드 12 참조.'
    }

# =====================================================================
# Slide 11: 조직 운영
# Table 0: 14 rows × 15 cols — header ['', '2012', ..., '2024', '평균']
#          rows: 12 departments + Total + 합계(F&B 제외)
# Table 1: 6 rows × 14 cols — header ['(단위)', '2012', ..., '2024']
#          rows: 급여+수당 / 종업원급여충당금 / 인건비 합 / 상승률 / 평균 인당 인건비
# =====================================================================
t = cells(11, 0)
if t:
    header = t[0]
    year_cols = []
    avg_col = None
    for i, h in enumerate(header):
        if re.match(r'^20\d{2}$', h.strip()):
            year_cols.append((i, h.strip()))
        elif '평균' in h.strip():
            avg_col = i
    years = [y for _, y in year_cols]
    departments = {}
    for row in t[1:]:
        label = row[0].strip()
        if not label: continue
        rec = {y: to_num(row[i]) for i, y in year_cols if i < len(row)}
        if avg_col is not None and avg_col < len(row):
            rec['avg'] = to_num(row[avg_col])
        departments[label] = rec
    org_headcount = {'years': years, 'departments': departments}
else:
    org_headcount = None

t = cells(11, 1)
if t:
    header = t[0]
    year_cols = [(i, h.strip()) for i, h in enumerate(header) if re.match(r'^20\d{2}$', h.strip())]
    years = [y for _, y in year_cols]
    labor_lines = {}
    for row in t[1:]:
        label = row[0].strip()
        if not label: continue
        labor_lines[label] = {y: to_num(row[i]) for i, y in year_cols if i < len(row)}
    labor_cost = {'years': years, 'lines': labor_lines, 'unit': 'IDR Million / persons'}
else:
    labor_cost = None

if org_headcount or labor_cost:
    out['organization_13y'] = {
        'source_slide': 11,
        'source_title': '10. OPEX - 조직 운영',
        'headcount': org_headcount,
        'labor_cost': labor_cost,
        'note': 'Headcount FY2012-FY2024 (13y). Avg ~240-250 total, ~170 ex-F&B.'
    }

# =====================================================================
# Slide 12: PBB (토지/빌딩세)
# Tables have complex merged-cell headers — store raw cells for downstream
# manual mapping; also try to extract a clean (year -> PBB total) series
# from rows 1-3 which contain 면적/공시가격/PBB.
# =====================================================================
t0 = cells(12, 0)
t1 = cells(12, 1)
pbb_raw = {}
if t0: pbb_raw['table_fy12_17'] = t0
if t1: pbb_raw['table_fy18_24'] = t1

# Attempt to derive a yearly PBB total (last numeric row labelled approximately "PBB" / "합계")
def extract_pbb_yearly(t, years):
    """Header pattern in slide 12 is non-standard. Years appear in header row 0
    at irregular positions. Sub-header row 1 has '공시가격' / 'PBB' labels per year.
    We locate the 'PBB' columns in row 1 and map them to years from row 0.
    """
    if not t or len(t) < 3: return {}
    h0 = t[0]
    h1 = t[1]
    yearly_pbb = {}
    # Walk row 1 and find PBB cols
    for ci, sub in enumerate(h1):
        if 'PBB' in sub:
            # Year is the nearest preceding year header in h0
            yr = None
            for j in range(ci, -1, -1):
                if j < len(h0) and re.match(r'^20\d{2}$', h0[j].strip()):
                    yr = h0[j].strip(); break
            if yr is None:
                # also check row 0 at same position
                continue
            # Sum the PBB column across all numeric rows (typically the last row holds total)
            colvals = []
            for r in t[2:]:
                if ci < len(r):
                    n = to_num(r[ci])
                    if n is not None: colvals.append(n)
            if colvals:
                yearly_pbb[yr] = colvals[-1]  # use last row as the total
    return yearly_pbb

pbb_yearly = {}
pbb_yearly.update(extract_pbb_yearly(t0, [str(y) for y in range(2012, 2018)]))
pbb_yearly.update(extract_pbb_yearly(t1, [str(y) for y in range(2018, 2025)]))

out['pbb_13y'] = {
    'source_slide': 12,
    'source_title': '11. OPEX – 토지/빌딩세',
    'unit': 'IDR (raw rupiah, NOT million) — verified against opex_14y["세금 및 법률 관련비용"] FY24=24,206 IDR Mil ≈ PBB FY24=23.3B IDR.',
    'note': 'PBB(토지·빌딩세) CAGR(2012~2024) ~17%. 전체 OPEX 약 30% 차지.',
    'yearly_pbb_total_idr': pbb_yearly,
    'raw_tables': pbb_raw,
    'caveat': '병합 셀이 많은 표 구조로 FY2013은 PBB 컬럼이 표 헤더에 별도 존재하지 않아 누락(필요 시 interpolation 권장). raw_tables는 검증용으로 보존.'
}

# =====================================================================
# Slide 13: 계약 + 물 사용
# Table 0: 22 rows × 3 cols — 21 contract rows
# Table 1: 5 rows × 7 cols — 4 water-usage lines × FY19-24
# =====================================================================
t = cells(13, 0)
contracts_list = []
if t and len(t) > 1:
    for row in t[1:]:
        if len(row) < 3 or not row[0].strip(): continue
        contracts_list.append({
            'counterparty': row[0].strip(),
            'contract': row[1].strip(),
            'revenue_classification': row[2].strip(),
        })

t = cells(13, 1)
water_lines = {}
water_years = []
if t and len(t) > 1:
    header = t[0]
    year_cols = [(i, h.strip()) for i, h in enumerate(header) if re.match(r'^20\d{2}$', h.strip())]
    water_years = [y for _, y in year_cols]
    for row in t[1:]:
        label = row[0].strip()
        if not label: continue
        water_lines[label] = {y: to_num(row[i]) for i, y in year_cols if i < len(row)}

out['contracts'] = {
    'source_slide': 13,
    'source_title': '12. 기타 - 주요 계약 목록 / 물 사용',
    'agreements': contracts_list,
    'agreements_count': len(contracts_list),
    'water_usage': {'years': water_years, 'lines': water_lines, 'unit': 'm³'} if water_lines else None,
}

# Update metadata
out['extracted'] = '2026-05-13 (v2: +capex_ppe/opex/org/pbb/contracts)'
out.setdefault('extraction_versions', {})['v2_added_at'] = '2026-05-13'
out['extraction_versions']['v2_sections'] = [
    'capex_ppe_14y', 'capex_auxiliary', 'opex_14y',
    'organization_13y', 'pbb_13y', 'contracts'
]

with open(OUT, 'w', encoding='utf-8') as f:
    json.dump(out, f, ensure_ascii=False, indent=2)

# Summary
print(f'Wrote {OUT}')
print(f'  capex_ppe_14y categories: {len(out.get("capex_ppe_14y",{}).get("categories",{}))}')
print(f'  opex_14y lines: {len(out.get("opex_14y",{}).get("lines",[]))} (+ total)')
print(f'  organization_13y departments: {len(out.get("organization_13y",{}).get("headcount",{}).get("departments",{}))}')
print(f'  organization_13y labor lines: {len(out.get("organization_13y",{}).get("labor_cost",{}).get("lines",{}))}')
print(f'  pbb_13y yearly entries: {len(out.get("pbb_13y",{}).get("yearly_pbb_total_idr",{}))}')
print(f'  contracts: {out.get("contracts",{}).get("agreements_count",0)} agreements')
